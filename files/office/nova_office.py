"""
💼 NOVA OFFICE — Generador de documentos profesionales con IA
==============================================================
Crea con un click:
  📄 Documentos Word (informes, cartas, ensayos, manuales)
  📊 Hojas Excel (presupuestos, planificadores, listas, tablas)
  🎯 Presentaciones PowerPoint (con estructura automática)
  📋 Currículums profesionales

Parte del ecosistema NOVA AI.
"""
import os, sys, json, threading, re
import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog
from datetime import datetime

APP_DIR = os.path.dirname(os.path.abspath(__file__))

# ============================================================
# VERSION (para el sistema de actualizaciones)
# ============================================================
VERSION = "1.0.0"
MODULO_ID = "office"

PARENT = os.path.dirname(APP_DIR)
CONFIG = os.path.join(PARENT, "config.json")
DOCS_DIR = os.path.join(APP_DIR, "documentos")
os.makedirs(DOCS_DIR, exist_ok=True)

sys.path.insert(0, APP_DIR)
from nova_ia_lib import llamar_llm, dialogo_config
try:
    from nova_brain import Brain
    BRAIN = Brain()
except Exception:
    BRAIN = None

# Dependencias opcionales
try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGBColor
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ============================================================
# PALETA
# ============================================================
C = {
    "bg":      "#0f172a",
    "bg2":     "#1e293b",
    "bg3":     "#334155",
    "fg":      "#f8fafc",
    "fg2":     "#94a3b8",
    "accent":  "#f59e0b",   # ámbar (Office Word)
    "excel":   "#22c55e",   # verde (Excel)
    "ppt":     "#ef4444",   # rojo (PowerPoint)
    "word":    "#3b82f6",   # azul (Word)
    "ok":      "#22c55e",
    "warn":    "#fbbf24",
    "err":     "#ef4444",
}

# ============================================================
# PLANTILLAS PREDEFINIDAS
# ============================================================
PLANTILLAS_WORD = {
    "📝 Ensayo académico":
        "ensayo académico formal de unas 800-1200 palabras sobre el tema. "
        "Incluye introducción, 3 argumentos principales con subtítulos, conclusión. "
        "Tono académico y referencias internas.",
    "📰 Artículo periodístico":
        "artículo periodístico con titular llamativo, entradilla en negrita, "
        "5-7 párrafos informativos y cita final.",
    "💌 Carta formal":
        "carta formal con encabezado (lugar y fecha), saludo apropiado, "
        "cuerpo claro de 3-4 párrafos y despedida formal.",
    "📋 Informe técnico":
        "informe técnico con: resumen ejecutivo, objetivos, metodología, "
        "resultados, conclusiones y recomendaciones. Usa subtítulos numerados.",
    "📖 Cuento corto":
        "cuento corto creativo con introducción que enganche, nudo con conflicto "
        "y desenlace sorprendente. Unas 600-1000 palabras.",
    "🎓 Trabajo escolar":
        "trabajo escolar bien estructurado con portada (al inicio en una línea), "
        "índice, 4-5 secciones temáticas y bibliografía sugerida.",
    "🏢 Memorándum empresarial":
        "memorándum empresarial formal con: De/Para/Fecha/Asunto al inicio, "
        "exposición clara del asunto en 2-3 párrafos, y conclusión accionable.",
    "📑 Manual de usuario":
        "manual de usuario con: introducción, requisitos, instalación paso a paso, "
        "guía de uso, solución de problemas y FAQ.",
}

PLANTILLAS_EXCEL = {
    "💰 Presupuesto mensual personal":
        "presupuesto mensual con categorías típicas: Vivienda, Alimentación, "
        "Transporte, Ocio, Ahorro, Otros. Para cada una un importe y un porcentaje. "
        "Datos realistas de ejemplo.",
    "📅 Planificador semanal":
        "planificador semanal: columnas = días de la semana (Lun-Dom), "
        "filas = horas (8:00 a 22:00 cada hora). Rellena tareas típicas de ejemplo.",
    "📋 Lista de tareas (To-Do)":
        "lista de tareas con columnas: ID, Tarea, Prioridad (Alta/Media/Baja), "
        "Estado (Pendiente/En curso/Hecho), Fecha límite, Notas. 15 tareas de ejemplo.",
    "📊 Tabla comparativa":
        "tabla comparativa de productos/opciones con columnas para Nombre, Precio, "
        "5 características distintas y Puntuación final.",
    "💼 Registro de gastos":
        "registro de gastos con columnas: Fecha, Descripción, Categoría, "
        "Importe (€), Método de pago. 20 movimientos de ejemplo.",
    "🎯 Plan de objetivos SMART":
        "plan de objetivos con columnas: Objetivo, Específico, Medible, Alcanzable, "
        "Relevante, Temporal (fecha), Estado. 8 objetivos de ejemplo.",
    "📈 Tracker de hábitos":
        "tracker mensual de hábitos: columnas = días (1-31), filas = hábitos "
        "(beber agua, ejercicio, lectura, meditar, dormir 8h). Marca X de ejemplo.",
    "🛒 Lista de la compra":
        "lista de la compra organizada por categoría (Frutas, Lácteos, Carnes, "
        "Limpieza, etc.) con columnas: Producto, Cantidad, Precio aprox.",
}

PLANTILLAS_PPT = {
    "🎤 Presentación corporativa":
        "presentación corporativa de 8 diapositivas: portada, problema, solución, "
        "producto, mercado, equipo, números clave, llamada a la acción.",
    "🎓 Presentación académica":
        "presentación académica de 10 diapositivas: portada, índice, introducción, "
        "objetivos, marco teórico, metodología, resultados, discusión, "
        "conclusiones, referencias.",
    "📖 Storytelling":
        "presentación narrativa de 6 diapositivas que cuenta una historia: "
        "personaje, mundo, conflicto, viaje, transformación, mensaje final.",
    "🚀 Pitch de startup":
        "pitch de startup de 10 diapositivas: problema, solución, demo, mercado, "
        "competencia, modelo de negocio, tracción, equipo, financiación, contacto.",
    "📊 Informe de resultados":
        "informe de resultados de 7 diapositivas: portada, resumen ejecutivo, "
        "KPIs, gráficos clave (descripción), análisis, conclusiones, próximos pasos.",
    "🎉 Presentación de evento":
        "presentación de un evento de 6 diapositivas: portada, ¿qué?, ¿cuándo?, "
        "¿dónde?, programa, registro/contacto.",
    "📚 Lección/Tutorial":
        "lección educativa de 8 diapositivas: portada, objetivos, qué aprenderás, "
        "concepto 1, concepto 2, concepto 3, ejercicio práctico, resumen.",
}

# ============================================================
# APP PRINCIPAL
# ============================================================
class NovaOffice:
    def __init__(self, root):
        self.root = root
        self.root.title("💼 NOVA Office — Documentos con IA")
        self.root.geometry("1180x780")
        self.root.configure(bg=C["bg"])

        self.cfg = self._cargar_config()
        self.proveedor = self.cfg.get("office_proveedor", "Groq")
        self.modelo = self.cfg.get("office_modelo", "llama-3.3-70b-versatile")
        self.api_key = self.cfg.get("api_keys", {}).get(self.proveedor, "")
        self.procesando = False

        self._construir_ui()

    def _cargar_config(self):
        try:
            with open(CONFIG, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return {"api_keys": {}}

    def _guardar_config(self):
        try:
            with open(CONFIG, "w", encoding="utf-8") as f:
                json.dump(self.cfg, f, indent=2)
        except Exception as e:
            print(f"Config: {e}")

    # ----------------------------------------------------------
    def _construir_ui(self):
        # HEADER
        h = tk.Frame(self.root, bg=C["bg2"], height=64)
        h.pack(fill="x"); h.pack_propagate(False)
        tk.Label(h, text="💼 NOVA Office",
                 bg=C["bg2"], fg=C["accent"],
                 font=("Segoe UI", 18, "bold")).pack(side="left", padx=20)
        tk.Label(h, text="Genera documentos profesionales con IA",
                 bg=C["bg2"], fg=C["fg2"],
                 font=("Segoe UI", 10, "italic")).pack(side="left", padx=5)

        for txt, cmd in [
            ("⚙ Config IA", self._config_ia),
            ("📂 Mis documentos", self._abrir_carpeta_docs),
            ("📖 Ayuda", self._mostrar_ayuda),
        ]:
            tk.Button(h, text=txt, command=cmd,
                bg=C["bg3"], fg=C["fg"], relief="flat",
                font=("Segoe UI", 9, "bold"), cursor="hand2",
                padx=12, pady=4).pack(side="right", padx=4, pady=10)

        # Notebook
        style = ttk.Style()
        try: style.theme_use("clam")
        except Exception: pass
        style.configure("TNotebook", background=C["bg"], borderwidth=0)
        style.configure("TNotebook.Tab", background=C["bg2"], foreground=C["fg"],
                         padding=[24, 10], font=("Segoe UI", 11, "bold"))
        style.map("TNotebook.Tab",
                   background=[("selected", C["accent"])],
                   foreground=[("selected", C["bg"])])

        nb = ttk.Notebook(self.root)
        nb.pack(fill="both", expand=True, padx=12, pady=12)

        self._tab_word(nb)
        self._tab_excel(nb)
        self._tab_ppt(nb)
        self._tab_cv(nb)

        # Status bar
        self.status = tk.Label(self.root,
            text="💼 NOVA Office listo. Elige Word / Excel / PowerPoint / CV.",
            bg=C["bg2"], fg=C["fg2"],
            font=("Segoe UI", 9), anchor="w", padx=14)
        self.status.pack(side="bottom", fill="x")

    # ============================================================
    # WORD
    # ============================================================
    def _tab_word(self, parent):
        tab = tk.Frame(parent, bg=C["bg"])
        parent.add(tab, text="📄 Word (.docx)")

        if not DOCX_OK:
            self._sin_libreria(tab, "python-docx", "pip install python-docx")
            return

        tk.Label(tab, text="📄 Generar documento Word con IA",
            bg=C["bg"], fg=C["word"],
            font=("Segoe UI", 13, "bold")).pack(anchor="w", padx=14, pady=12)

        tk.Label(tab, text="¿Qué tipo de documento?",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14)
        self.word_tipo = tk.StringVar(value=list(PLANTILLAS_WORD.keys())[0])
        ttk.Combobox(tab, textvariable=self.word_tipo,
            values=list(PLANTILLAS_WORD.keys()),
            state="readonly", font=("Segoe UI", 10)).pack(
            fill="x", padx=14, pady=4, ipady=2)

        tk.Label(tab, text="Tema/asunto:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14, pady=(10, 2))
        self.word_tema = tk.Entry(tab, bg=C["bg2"], fg=C["fg"],
            insertbackground=C["fg"], relief="flat", font=("Segoe UI", 11))
        self.word_tema.pack(fill="x", padx=14, ipady=8)

        tk.Label(tab, text="Instrucciones adicionales (opcional):",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14, pady=(10, 2))
        self.word_extra = tk.Text(tab, bg=C["bg2"], fg=C["fg"],
            insertbackground=C["fg"], relief="flat", font=("Segoe UI", 10),
            height=3, wrap="word")
        self.word_extra.pack(fill="x", padx=14, ipady=4)

        tk.Button(tab, text="✨ Generar documento Word",
            command=self._generar_word,
            bg=C["word"], fg="white", relief="flat",
            font=("Segoe UI", 11, "bold"), cursor="hand2",
            pady=10).pack(fill="x", padx=14, pady=14)

        # Vista previa
        tk.Label(tab, text="Vista previa:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14)
        self.word_preview = scrolledtext.ScrolledText(tab,
            bg=C["bg2"], fg=C["fg"], font=("Segoe UI", 10),
            wrap="word", relief="flat", padx=12, pady=10)
        self.word_preview.pack(fill="both", expand=True, padx=14, pady=(4, 14))

    def _generar_word(self):
        tema = self.word_tema.get().strip()
        if not tema:
            messagebox.showwarning("⚠️", "Escribe un tema."); return
        if not self._check_api(): return
        if self.procesando: return

        tipo = self.word_tipo.get()
        plantilla = PLANTILLAS_WORD[tipo]
        extra = self.word_extra.get("1.0", "end").strip()

        self.procesando = True
        self.word_preview.delete("1.0", "end")
        self.word_preview.insert("1.0", "⏳ Generando documento... (10-40s)\n")
        self.status.config(text=f"⏳ Generando Word: {tipo}...", fg=C["accent"])

        def tarea():
            system = (
                "Eres un asistente experto en redacción. "
                "Generas documentos en español, bien estructurados. "
                "Usa este formato MARKDOWN:\n"
                "  # Título principal\n"
                "  ## Subtítulos\n"
                "  ### Sub-subtítulos\n"
                "  **negrita** y texto normal.\n"
                "  - listas con guiones\n"
                "  1. listas numeradas\n"
                "NO añadas comentarios extra fuera del documento."
            )
            user = f"""Crea un {plantilla}

Tema: "{tema}"
{("Notas adicionales: " + extra) if extra else ""}

Devuelve el documento completo en MARKDOWN. Asegúrate de:
- Empezar con # (título principal)
- Usar ## para secciones grandes
- Usar **negrita** para destacar
- Listas cuando sea apropiado"""

            resp, err = llamar_llm(self.proveedor, self.modelo, self.api_key,
                                    system, user)
            self.root.after(0, self._word_completo, resp, err, tema, tipo)

        threading.Thread(target=tarea, daemon=True).start()

    def _word_completo(self, resp, err, tema, tipo):
        self.procesando = False
        if err:
            self.word_preview.delete("1.0", "end")
            self.word_preview.insert("1.0", f"❌ Error: {err}")
            self.status.config(text=f"❌ {err}", fg=C["err"])
            return

        # Mostrar preview
        self.word_preview.delete("1.0", "end")
        self.word_preview.insert("1.0", resp)

        # Guardar Word
        try:
            nombre = self._slug(tema) + "_" + datetime.now().strftime("%Y%m%d_%H%M")
            ruta = os.path.join(DOCS_DIR, f"{nombre}.docx")
            self._markdown_a_docx(resp, ruta, titulo_default=tema)
            # Registrar en el Brain global
            if BRAIN:
                BRAIN.evento("documento_creado", "office",
                    {"formato": "word", "tipo": tipo, "titulo": tema, "ruta": ruta})
                BRAIN.stat("office", "docs_word", incrementar=1)
            self.status.config(
                text=f"✅ Word guardado: {os.path.basename(ruta)}",
                fg=C["ok"])
            if messagebox.askyesno("✅ Documento listo",
                f"Word guardado:\n{ruta}\n\n¿Abrirlo ahora?"):
                self._abrir(ruta)
        except Exception as e:
            self.status.config(text=f"❌ Error al guardar: {e}", fg=C["err"])

    def _markdown_a_docx(self, md, ruta, titulo_default="Documento"):
        """Convierte markdown básico a Word con estilos."""
        doc = Document()

        # Configurar estilos
        styles = doc.styles
        try:
            normal = styles["Normal"]
            normal.font.name = "Calibri"
            normal.font.size = Pt(11)
        except Exception:
            pass

        lineas = md.split("\n")
        for linea in lineas:
            s = linea.rstrip()
            if not s:
                doc.add_paragraph()
                continue
            # Encabezados
            if s.startswith("# "):
                p = doc.add_heading(s[2:].strip(), level=0)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif s.startswith("## "):
                doc.add_heading(s[3:].strip(), level=1)
            elif s.startswith("### "):
                doc.add_heading(s[4:].strip(), level=2)
            elif s.startswith("#### "):
                doc.add_heading(s[5:].strip(), level=3)
            # Listas
            elif s.lstrip().startswith(("- ", "* ", "+ ")):
                texto = s.lstrip()[2:].strip()
                p = doc.add_paragraph(style="List Bullet")
                self._add_runs(p, texto)
            elif re.match(r"^\s*\d+\.\s", s):
                texto = re.sub(r"^\s*\d+\.\s", "", s)
                p = doc.add_paragraph(style="List Number")
                self._add_runs(p, texto)
            # Línea separadora
            elif s.strip() in ("---", "***"):
                doc.add_paragraph("─" * 40).alignment = WD_ALIGN_PARAGRAPH.CENTER
            # Cita
            elif s.startswith("> "):
                p = doc.add_paragraph(s[2:].strip(), style="Intense Quote")
            # Párrafo normal
            else:
                p = doc.add_paragraph()
                self._add_runs(p, s)

        doc.save(ruta)

    def _add_runs(self, paragraph, texto):
        """Añade runs respetando **negrita** y *cursiva*."""
        # Negrita primero
        partes = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", texto)
        for parte in partes:
            if not parte: continue
            if parte.startswith("**") and parte.endswith("**"):
                r = paragraph.add_run(parte[2:-2])
                r.bold = True
            elif parte.startswith("*") and parte.endswith("*") and len(parte) > 2:
                r = paragraph.add_run(parte[1:-1])
                r.italic = True
            else:
                paragraph.add_run(parte)

    # ============================================================
    # EXCEL
    # ============================================================
    def _tab_excel(self, parent):
        tab = tk.Frame(parent, bg=C["bg"])
        parent.add(tab, text="📊 Excel (.xlsx)")

        if not XLSX_OK:
            self._sin_libreria(tab, "openpyxl", "pip install openpyxl")
            return

        tk.Label(tab, text="📊 Generar hoja Excel con IA",
            bg=C["bg"], fg=C["excel"],
            font=("Segoe UI", 13, "bold")).pack(anchor="w", padx=14, pady=12)

        tk.Label(tab, text="¿Qué tipo de hoja?",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14)
        self.excel_tipo = tk.StringVar(value=list(PLANTILLAS_EXCEL.keys())[0])
        ttk.Combobox(tab, textvariable=self.excel_tipo,
            values=list(PLANTILLAS_EXCEL.keys()),
            state="readonly", font=("Segoe UI", 10)).pack(
            fill="x", padx=14, pady=4, ipady=2)

        tk.Label(tab, text="Detalles o personalización (opcional):",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14, pady=(10, 2))
        self.excel_extra = tk.Text(tab, bg=C["bg2"], fg=C["fg"],
            insertbackground=C["fg"], relief="flat", font=("Segoe UI", 10),
            height=3, wrap="word")
        self.excel_extra.pack(fill="x", padx=14, ipady=4)

        tk.Button(tab, text="✨ Generar Excel",
            command=self._generar_excel,
            bg=C["excel"], fg="white", relief="flat",
            font=("Segoe UI", 11, "bold"), cursor="hand2",
            pady=10).pack(fill="x", padx=14, pady=14)

        tk.Label(tab, text="Vista previa CSV:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14)
        self.excel_preview = scrolledtext.ScrolledText(tab,
            bg=C["bg2"], fg=C["fg"], font=("Consolas", 10),
            wrap="none", relief="flat", padx=12, pady=10)
        self.excel_preview.pack(fill="both", expand=True, padx=14, pady=(4, 14))

    def _generar_excel(self):
        if not self._check_api(): return
        if self.procesando: return
        tipo = self.excel_tipo.get()
        plantilla = PLANTILLAS_EXCEL[tipo]
        extra = self.excel_extra.get("1.0", "end").strip()

        self.procesando = True
        self.excel_preview.delete("1.0", "end")
        self.excel_preview.insert("1.0", "⏳ Generando datos...\n")
        self.status.config(text=f"⏳ Generando Excel: {tipo}...", fg=C["accent"])

        def tarea():
            system = (
                "Eres experto en hojas de cálculo. Generas datos en formato CSV "
                "(separados por |) listos para Excel. La primera línea son los "
                "ENCABEZADOS de columnas. Los datos son realistas y útiles. "
                "NO añadas explicaciones, SÓLO el CSV con separador '|'."
            )
            user = f"""Genera datos para: {plantilla}

{("Especificaciones extra: " + extra) if extra else ""}

Devuelve SOLO el CSV usando '|' como separador. Ejemplo:
Encabezado1|Encabezado2|Encabezado3
Dato1|Dato2|Dato3
DatoA|DatoB|DatoC

(SIN explicaciones, SIN backticks, SIN texto adicional)"""

            resp, err = llamar_llm(self.proveedor, self.modelo, self.api_key,
                                    system, user)
            self.root.after(0, self._excel_completo, resp, err, tipo)

        threading.Thread(target=tarea, daemon=True).start()

    def _excel_completo(self, resp, err, tipo):
        self.procesando = False
        if err:
            self.excel_preview.delete("1.0", "end")
            self.excel_preview.insert("1.0", f"❌ Error: {err}")
            self.status.config(text=f"❌ {err}", fg=C["err"])
            return

        # Limpiar backticks si llegan
        resp = re.sub(r"```\w*\n?", "", resp).strip("`\n ")

        self.excel_preview.delete("1.0", "end")
        self.excel_preview.insert("1.0", resp)

        # Parsear y guardar Excel
        try:
            filas = [l for l in resp.split("\n") if "|" in l]
            if not filas:
                self.status.config(text="❌ La IA no devolvió datos válidos",
                                    fg=C["err"])
                return
            data = [[c.strip() for c in f.split("|")] for f in filas]

            wb = Workbook()
            ws = wb.active
            ws.title = tipo.replace(" ", "_")[:30]

            # Estilos
            header_fill = PatternFill(start_color="22C55E",
                                       end_color="22C55E",
                                       fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=11)
            cell_font = Font(size=10)
            thin = Side(border_style="thin", color="CCCCCC")
            border = Border(left=thin, right=thin, top=thin, bottom=thin)

            for r_idx, fila in enumerate(data, 1):
                for c_idx, valor in enumerate(fila, 1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=valor)
                    cell.border = border
                    if r_idx == 1:
                        cell.fill = header_fill
                        cell.font = header_font
                        cell.alignment = Alignment(horizontal="center",
                                                    vertical="center")
                    else:
                        cell.font = cell_font

            # Ajustar anchos
            for col in range(1, len(data[0]) + 1):
                ws.column_dimensions[get_column_letter(col)].width = 18
            ws.row_dimensions[1].height = 28

            nombre = self._slug(tipo) + "_" + datetime.now().strftime("%Y%m%d_%H%M")
            ruta = os.path.join(DOCS_DIR, f"{nombre}.xlsx")
            wb.save(ruta)
            self.status.config(
                text=f"✅ Excel guardado: {os.path.basename(ruta)}", fg=C["ok"])

            if messagebox.askyesno("✅ Excel listo",
                f"Excel guardado:\n{ruta}\n\n¿Abrirlo?"):
                self._abrir(ruta)
        except Exception as e:
            self.status.config(text=f"❌ Error: {e}", fg=C["err"])

    # ============================================================
    # POWERPOINT
    # ============================================================
    def _tab_ppt(self, parent):
        tab = tk.Frame(parent, bg=C["bg"])
        parent.add(tab, text="🎯 PowerPoint (.pptx)")

        if not PPTX_OK:
            self._sin_libreria(tab, "python-pptx", "pip install python-pptx")
            return

        tk.Label(tab, text="🎯 Generar presentación PowerPoint con IA",
            bg=C["bg"], fg=C["ppt"],
            font=("Segoe UI", 13, "bold")).pack(anchor="w", padx=14, pady=12)

        tk.Label(tab, text="Tipo de presentación:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14)
        self.ppt_tipo = tk.StringVar(value=list(PLANTILLAS_PPT.keys())[0])
        ttk.Combobox(tab, textvariable=self.ppt_tipo,
            values=list(PLANTILLAS_PPT.keys()),
            state="readonly", font=("Segoe UI", 10)).pack(
            fill="x", padx=14, pady=4, ipady=2)

        tk.Label(tab, text="Tema/asunto:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14, pady=(10, 2))
        self.ppt_tema = tk.Entry(tab, bg=C["bg2"], fg=C["fg"],
            insertbackground=C["fg"], relief="flat", font=("Segoe UI", 11))
        self.ppt_tema.pack(fill="x", padx=14, ipady=8)

        # Tema visual
        tema_frame = tk.Frame(tab, bg=C["bg"])
        tema_frame.pack(fill="x", padx=14, pady=(10, 0))
        tk.Label(tema_frame, text="Color del tema:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(side="left")
        self.ppt_color = tk.StringVar(value="Azul corporativo")
        ttk.Combobox(tema_frame, textvariable=self.ppt_color,
            values=["Azul corporativo", "Rojo elegante", "Verde natural",
                    "Morado creativo", "Naranja energético", "Gris ejecutivo"],
            state="readonly", width=22, font=("Segoe UI", 10)).pack(
            side="left", padx=8)

        tk.Button(tab, text="✨ Generar PowerPoint",
            command=self._generar_ppt,
            bg=C["ppt"], fg="white", relief="flat",
            font=("Segoe UI", 11, "bold"), cursor="hand2",
            pady=10).pack(fill="x", padx=14, pady=14)

        tk.Label(tab, text="Vista previa de las diapositivas:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14)
        self.ppt_preview = scrolledtext.ScrolledText(tab,
            bg=C["bg2"], fg=C["fg"], font=("Segoe UI", 10),
            wrap="word", relief="flat", padx=12, pady=10)
        self.ppt_preview.pack(fill="both", expand=True, padx=14, pady=(4, 14))

    def _generar_ppt(self):
        tema = self.ppt_tema.get().strip()
        if not tema:
            messagebox.showwarning("⚠️", "Escribe un tema."); return
        if not self._check_api(): return
        if self.procesando: return

        tipo = self.ppt_tipo.get()
        plantilla = PLANTILLAS_PPT[tipo]
        color = self.ppt_color.get()

        self.procesando = True
        self.ppt_preview.delete("1.0", "end")
        self.ppt_preview.insert("1.0", "⏳ Generando diapositivas... (10-40s)\n")
        self.status.config(text=f"⏳ Generando PowerPoint...", fg=C["accent"])

        def tarea():
            system = (
                "Eres un experto en presentaciones. Generas contenido para "
                "diapositivas usando este formato EXACTO:\n\n"
                "DIAPOSITIVA 1:\n"
                "TÍTULO: [título corto, max 8 palabras]\n"
                "CONTENIDO:\n"
                "- Punto 1 (corto y claro)\n"
                "- Punto 2\n"
                "- Punto 3\n"
                "\n"
                "DIAPOSITIVA 2:\n"
                "...\n\n"
                "Reglas:\n"
                "- Máximo 5 puntos por diapositiva\n"
                "- Cada punto MUY conciso (máx 12 palabras)\n"
                "- NO añadas notas ni comentarios fuera del formato\n"
                "- Lenguaje claro y profesional"
            )
            user = f"""Crea una {plantilla}

Tema: "{tema}"

Devuelve exactamente las diapositivas en el formato indicado."""

            resp, err = llamar_llm(self.proveedor, self.modelo, self.api_key,
                                    system, user)
            self.root.after(0, self._ppt_completo, resp, err, tema, tipo, color)

        threading.Thread(target=tarea, daemon=True).start()

    def _ppt_completo(self, resp, err, tema, tipo, color):
        self.procesando = False
        if err:
            self.ppt_preview.delete("1.0", "end")
            self.ppt_preview.insert("1.0", f"❌ Error: {err}")
            self.status.config(text=f"❌ {err}", fg=C["err"])
            return

        self.ppt_preview.delete("1.0", "end")
        self.ppt_preview.insert("1.0", resp)

        # Parsear
        try:
            slides = self._parsear_ppt(resp)
            if not slides:
                self.status.config(
                    text="❌ Formato no reconocido", fg=C["err"])
                return

            nombre = self._slug(tema) + "_" + datetime.now().strftime("%Y%m%d_%H%M")
            ruta = os.path.join(DOCS_DIR, f"{nombre}.pptx")
            self._crear_pptx(slides, ruta, color)

            self.status.config(
                text=f"✅ PowerPoint guardado: {os.path.basename(ruta)} ({len(slides)} slides)",
                fg=C["ok"])
            if messagebox.askyesno("✅ Listo",
                f"PowerPoint guardado:\n{ruta}\n\n¿Abrirlo?"):
                self._abrir(ruta)
        except Exception as e:
            self.status.config(text=f"❌ Error: {e}", fg=C["err"])

    def _parsear_ppt(self, texto):
        slides = []
        bloques = re.split(r"DIAPOSITIVA\s+\d+\s*:", texto, flags=re.IGNORECASE)
        for b in bloques[1:]:
            titulo_m = re.search(r"T[ÍI]TULO\s*:\s*(.+?)(?=\n)",
                                  b, re.IGNORECASE)
            contenido_m = re.search(
                r"CONTENIDO\s*:\s*(.+?)(?=DIAPOSITIVA|\Z)",
                b, re.DOTALL | re.IGNORECASE)
            titulo = titulo_m.group(1).strip() if titulo_m else "Diapositiva"
            puntos = []
            if contenido_m:
                for linea in contenido_m.group(1).split("\n"):
                    s = linea.strip()
                    if s.startswith(("-", "*", "•")):
                        puntos.append(s.lstrip("-*• ").strip())
            slides.append({"titulo": titulo, "puntos": puntos})
        return slides

    def _crear_pptx(self, slides, ruta, color_tema):
        colores_temas = {
            "Azul corporativo": ("1E40AF", "DBEAFE", "1E3A8A"),
            "Rojo elegante": ("991B1B", "FECACA", "7F1D1D"),
            "Verde natural": ("166534", "DCFCE7", "14532D"),
            "Morado creativo": ("6B21A8", "F3E8FF", "581C87"),
            "Naranja energético": ("C2410C", "FFEDD5", "9A3412"),
            "Gris ejecutivo": ("374151", "F3F4F6", "1F2937"),
        }
        c_main, c_bg, c_acc = colores_temas.get(color_tema,
                                                  colores_temas["Azul corporativo"])

        prs = Presentation()
        prs.slide_width = PInches(13.33)
        prs.slide_height = PInches(7.5)

        for i, slide_data in enumerate(slides):
            layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(layout)

            # Fondo del título
            from pptx.shapes.autoshape import Shape
            from pptx.enum.shapes import MSO_SHAPE

            # Barra superior de color
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                PInches(0), PInches(0), PInches(13.33), PInches(1.2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = PRGBColor.from_string(c_main)
            bar.line.fill.background()

            # Título
            tx = slide.shapes.add_textbox(PInches(0.5), PInches(0.25),
                                            PInches(12.3), PInches(0.8))
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["titulo"]
            p.font.size = PPt(32)
            p.font.bold = True
            p.font.color.rgb = PRGBColor.from_string("FFFFFF")

            # Contenido
            cx = slide.shapes.add_textbox(PInches(0.7), PInches(1.7),
                                            PInches(12), PInches(5.5))
            cf = cx.text_frame
            cf.word_wrap = True

            for j, punto in enumerate(slide_data["puntos"]):
                if j == 0:
                    par = cf.paragraphs[0]
                else:
                    par = cf.add_paragraph()
                par.text = "• " + punto
                par.font.size = PPt(20)
                par.font.color.rgb = PRGBColor.from_string("1F2937")
                par.space_after = PPt(14)

            # Número de slide en esquina
            num_box = slide.shapes.add_textbox(PInches(12.5), PInches(7),
                                                 PInches(0.7), PInches(0.4))
            nf = num_box.text_frame
            np = nf.paragraphs[0]
            np.text = f"{i+1} / {len(slides)}"
            np.font.size = PPt(10)
            np.font.color.rgb = PRGBColor.from_string("9CA3AF")

        prs.save(ruta)

    # ============================================================
    # CURRÍCULUM (CV)
    # ============================================================
    def _tab_cv(self, parent):
        tab = tk.Frame(parent, bg=C["bg"])
        parent.add(tab, text="📋 CV/Currículum")

        if not DOCX_OK:
            self._sin_libreria(tab, "python-docx", "pip install python-docx")
            return

        tk.Label(tab, text="📋 Genera tu CV profesional con IA",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 13, "bold")).pack(anchor="w", padx=14, pady=12)

        tk.Label(tab,
            text="Cuéntame sobre ti y la IA armará tu CV.\n"
                 "Incluye: tus estudios, experiencia, habilidades, idiomas, "
                 "y al puesto al que aspiras.",
            bg=C["bg"], fg=C["fg2"],
            font=("Segoe UI", 9, "italic"),
            justify="left").pack(anchor="w", padx=14)

        tk.Label(tab, text="📝 Cuéntame sobre ti:",
            bg=C["bg"], fg=C["accent"],
            font=("Segoe UI", 10, "bold")).pack(anchor="w", padx=14, pady=(10, 2))
        self.cv_texto = scrolledtext.ScrolledText(tab,
            bg=C["bg2"], fg=C["fg"], insertbackground=C["fg"],
            font=("Segoe UI", 10), wrap="word", relief="flat",
            height=10, padx=12, pady=10)
        self.cv_texto.pack(fill="both", expand=True, padx=14, pady=4)
        self.cv_texto.insert("1.0",
            "Nombre: \n"
            "Email: \n"
            "Teléfono: \n"
            "Ciudad: \n\n"
            "Estudios: \n\n"
            "Experiencia: \n\n"
            "Habilidades: \n\n"
            "Idiomas: \n\n"
            "Puesto al que aspiro: \n")

        tk.Button(tab, text="✨ Generar mi CV profesional",
            command=self._generar_cv,
            bg=C["accent"], fg="white", relief="flat",
            font=("Segoe UI", 11, "bold"), cursor="hand2",
            pady=10).pack(fill="x", padx=14, pady=14)

    def _generar_cv(self):
        datos = self.cv_texto.get("1.0", "end").strip()
        if len(datos) < 50:
            messagebox.showwarning("⚠️", "Escribe más información sobre ti.")
            return
        if not self._check_api(): return
        if self.procesando: return

        self.procesando = True
        self.status.config(text="⏳ Generando tu CV profesional...",
                            fg=C["accent"])

        def tarea():
            system = (
                "Eres un experto en RRHH y redacción de currículums. "
                "Generas CVs profesionales, claros y atractivos en español. "
                "Usa formato MARKDOWN:\n"
                "  # NOMBRE COMPLETO\n"
                "  ## Datos de contacto / Sección...\n"
                "  - Items con guiones\n"
                "  **Negrita** para destacar puestos/empresas/títulos"
            )
            user = f"""Crea un CURRÍCULUM PROFESIONAL bien estructurado a partir de estos datos:

{datos}

Estructura recomendada:
# NOMBRE
## 📧 Contacto
## 🎯 Perfil profesional (2-3 frases que enganchen)
## 💼 Experiencia laboral
## 🎓 Formación académica
## 🌐 Idiomas
## 💡 Habilidades técnicas
## ⚡ Soft skills

Mejora la redacción de cada sección. Si falta info, no la inventes (deja indicado [completar])."""

            resp, err = llamar_llm(self.proveedor, self.modelo, self.api_key,
                                    system, user)
            self.root.after(0, self._cv_completo, resp, err)

        threading.Thread(target=tarea, daemon=True).start()

    def _cv_completo(self, resp, err):
        self.procesando = False
        if err:
            self.status.config(text=f"❌ {err}", fg=C["err"])
            messagebox.showerror("❌", err); return
        try:
            nombre = "CV_" + datetime.now().strftime("%Y%m%d_%H%M")
            ruta = os.path.join(DOCS_DIR, f"{nombre}.docx")
            self._markdown_a_docx(resp, ruta, titulo_default="Currículum")
            self.status.config(
                text=f"✅ CV guardado: {os.path.basename(ruta)}", fg=C["ok"])
            if messagebox.askyesno("✅ ¡CV listo!",
                f"Tu CV está en:\n{ruta}\n\n¿Abrirlo ahora?"):
                self._abrir(ruta)
        except Exception as e:
            self.status.config(text=f"❌ {e}", fg=C["err"])

    # ============================================================
    # UTILIDADES
    # ============================================================
    def _sin_libreria(self, tab, nombre, comando):
        tk.Label(tab,
            text=f"⚠️ Necesitas instalar '{nombre}'\n\n"
                 f"Abre cmd y ejecuta:\n  {comando}\n\n"
                 "Después reinicia NOVA Office.",
            bg=C["bg"], fg=C["warn"],
            font=("Consolas", 11),
            justify="center").pack(expand=True)

    def _check_api(self):
        if not self.api_key and self.proveedor != "Ollama":
            messagebox.showwarning("⚠️",
                "Configura tu API Key en ⚙ Config IA primero.")
            return False
        return True

    def _slug(self, texto):
        s = re.sub(r"[^\w\s-]", "", texto).strip().lower()
        s = re.sub(r"[-\s]+", "_", s)
        return s[:40] or "documento"

    def _abrir(self, ruta):
        try: os.startfile(ruta)
        except AttributeError:
            import subprocess
            subprocess.Popen(["xdg-open", ruta])

    def _abrir_carpeta_docs(self):
        self._abrir(DOCS_DIR)

    def _config_ia(self):
        cfg = {"proveedor": self.proveedor, "modelo": self.modelo,
               "api_key": self.api_key, "api_keys": self.cfg.get("api_keys", {})}
        col = {"bg": C["bg"], "bg2": C["bg2"], "fg": C["fg"],
               "accent": C["accent"], "warn": C["warn"], "ok": C["ok"]}
        def on_save(n):
            self.proveedor = n["proveedor"]; self.modelo = n["modelo"]
            self.api_key = n["api_key"]
            self.cfg["office_proveedor"] = self.proveedor
            self.cfg["office_modelo"] = self.modelo
            self.cfg["api_keys"] = n["api_keys"]
            self._guardar_config()
        dialogo_config(self.root, cfg, on_save,
            titulo="⚙ Configuración IA — Office", colores=col)

    def _mostrar_ayuda(self):
        ayuda = """💼 NOVA OFFICE — Guía rápida

📄 WORD
   • Elige el tipo (ensayo, informe, carta...)
   • Escribe el tema
   • La IA genera el documento completo
   • Se guarda como .docx en /documentos

📊 EXCEL
   • Elige el tipo (presupuesto, planificador...)
   • Personaliza si quieres
   • La IA genera datos realistas
   • Se guarda como .xlsx con estilo profesional

🎯 POWERPOINT
   • Elige el tipo de presentación
   • Tema + color
   • La IA crea diapositivas con título + viñetas
   • Se guarda como .pptx con tu color elegido

📋 CV
   • Rellena tus datos en el formulario
   • La IA mejora la redacción profesional
   • Se guarda como .docx listo para enviar

💾 Todos los archivos van a:
   {dir}

⚙ Config IA: elige proveedor (Groq/Gemini/Ollama...)
   y la lista de modelos se carga automáticamente.

💡 TRUCO: cuanto MÁS DETALLES des en "tema" o
   "instrucciones extra", mejor saldrá el resultado.""".format(dir=DOCS_DIR)

        dlg = tk.Toplevel(self.root)
        dlg.title("📖 Ayuda")
        dlg.geometry("620x520")
        dlg.configure(bg=C["bg"])
        txt = scrolledtext.ScrolledText(dlg, bg=C["bg2"], fg=C["fg"],
            font=("Consolas", 10), wrap="word", relief="flat",
            padx=12, pady=12)
        txt.pack(fill="both", expand=True, padx=12, pady=12)
        txt.insert("1.0", ayuda)
        txt.config(state="disabled")


if __name__ == "__main__":
    try:
        from nova_ui import splash_screen, Theme
        splash_screen("NOVA Office",
                       subtitulo="Word, Excel, PowerPoint y CV con IA",
                       color_acento=Theme.accent4,
                       tareas=[
                           ("Cargando plantillas...", None),
                           ("Preparando editor...", None),
                           ("Suite lista...", None),
                       ],
                       duracion_min=1.3)
    except Exception as e:
        print(f"Splash: {e}")

    root = tk.Tk()
    app = NovaOffice(root)
    root.mainloop()
